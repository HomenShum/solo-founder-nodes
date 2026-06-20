#!/usr/bin/env python3
"""Generate office sample sources (xlsx/pptx/docx) + a shared claims file to dogfood cite.py
across all artifact types. Same content as sample.pdf (see make_sample.py).
Run: python make_samples.py   (needs openpyxl, python-pptx, python-docx)
"""
import json
import os
import sys

for _s in (sys.stdout, sys.stderr):
    try:
        _s.reconfigure(encoding="utf-8")
    except Exception:
        pass

HERE = os.path.dirname(os.path.abspath(__file__))
TITLE = "Sample Co - FY2025 Financial Summary"
LINES = [
    "Total revenue for fiscal year 2025 was $4.2 million, up 38% year over year.",
    "The company operates two enterprise pilot programs as of Q4 2025.",
    "Cash on hand at year end was $1.1 million.",
    "Headcount grew to 12 full-time employees during the year.",
]

# XLSX — each statement in a cell of column A
import openpyxl
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Summary"
ws["A1"] = TITLE
for i, line in enumerate(LINES, start=3):
    ws.cell(row=i, column=1, value=line)
wb.save(os.path.join(HERE, "sample.xlsx"))

# PPTX — each statement in its own text box on a blank slide
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)).text_frame.text = TITLE
y = 1.3
for line in LINES:
    slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5)).text_frame.text = line
    y += 0.7
prs.save(os.path.join(HERE, "sample.pptx"))

# DOCX — each statement a paragraph
import docx
d = docx.Document()
d.add_heading(TITLE, level=1)
for line in LINES:
    d.add_paragraph(line)
d.save(os.path.join(HERE, "sample.docx"))

# Shared claims (work across xlsx/pptx/docx — quotes are substrings present in all three)
claims = [
    {"id": "c1", "claim": "FY2025 revenue was $4.2M",
     "quote": "Total revenue for fiscal year 2025 was $4.2 million"},
    {"id": "c2", "claim": "The company runs two pilots",
     "quote": "operates two enterprise pilot programs"},
    {"id": "c3", "claim": "The company was profitable in 2025",
     "quote": "net income was positive in 2025"},
]
with open(os.path.join(HERE, "claims_office.json"), "w", encoding="utf-8") as f:
    json.dump(claims, f, indent=2)

print("wrote sample.xlsx, sample.pptx, sample.docx, claims_office.json")
print("c1/c2 should be SUPPORTED in every format; c3 (net income) should be UNSUPPORTED.")
