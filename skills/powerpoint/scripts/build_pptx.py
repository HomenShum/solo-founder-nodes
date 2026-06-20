#!/usr/bin/env python3
"""build_pptx.py — render a deck_plan.json into a real .pptx.

Self-contained: needs only `python-pptx` (pip install python-pptx). It renders
the honesty model visibly — verified claims get a muted source footer,
needs_review claims get an amber "needs review" marker, and a closing
"To Verify Before Presenting" slide is appended automatically from every
needs_review claim in the deck.

Usage:
    python build_pptx.py deck_plan.json deck.pptx
"""

import argparse
import json
import sys

for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8")
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx is required. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(2)

DEFAULT_THEME = {
    "primary": "B85042",     # terracotta
    "bg_dark": "1E2024",     # near-black
    "bg_light": "F7F4EF",    # warm cream
    "text_dark": "20232A",
    "text_light": "F7F4EF",
    "accent": "C8772E",      # burnt orange
    "warn": "C9892F",        # amber for needs_review
    "muted": "8A8A82",
}

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def color(theme, key):
    return RGBColor.from_string(theme.get(key, DEFAULT_THEME[key]))


def fill_bg(slide, theme, key):
    """Full-bleed background rectangle (works across python-pptx versions)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color(theme, key)
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_box(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    return tf


def style_run(run, text, size, rgb, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb


def source_label(plan, source):
    """Resolve a source id from the registry to its human label; else return as-is."""
    if not source:
        return ""
    for s in plan.get("sources", []) or []:
        if s.get("id") == source:
            return s.get("label") or s.get("url") or source
    return source


# ---- layouts ---------------------------------------------------------------

def render_title(slide, plan, theme, s):
    fill_bg(slide, theme, "bg_dark")
    tf = add_box(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
    p = tf.paragraphs[0]
    style_run(p.add_run(), s.get("title", ""), 48, color(theme, "text_light"), bold=True)
    if s.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        style_run(p2.add_run(), s["subtitle"], 22, color(theme, "accent"))


def render_closing(slide, plan, theme, s):
    fill_bg(slide, theme, "bg_dark")
    tf = add_box(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.4), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    style_run(p.add_run(), s.get("title", ""), 40, color(theme, "text_light"), bold=True)
    if s.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        style_run(p2.add_run(), s["subtitle"], 20, color(theme, "accent"))


def render_content(slide, plan, theme, s):
    fill_bg(slide, theme, "bg_light")
    # Title
    tf = add_box(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
    style_run(tf.paragraphs[0].add_run(), s.get("title", ""), 34, color(theme, "primary"), bold=True)

    # Bullets
    body = add_box(slide, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.6))
    sources_used = []
    first = True
    for b in s.get("bullets", []) or []:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.space_after = Pt(12)
        style_run(p.add_run(), "•  ", 18, color(theme, "primary"), bold=True)
        style_run(p.add_run(), b.get("text", ""), 18, color(theme, "text_dark"))
        status = b.get("status")
        if status == "needs_review":
            style_run(p.add_run(), "   ⚠ needs review", 13, color(theme, "warn"), bold=True)
        elif status == "manual":
            style_run(p.add_run(), "   (your note)", 12, color(theme, "muted"), italic=True)
        elif status == "verified":
            lbl = source_label(plan, b.get("source"))
            if lbl:
                sources_used.append(lbl)

    # Citation footer for verified claims
    if sources_used:
        foot = add_box(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
        uniq = list(dict.fromkeys(sources_used))
        style_run(foot.paragraphs[0].add_run(), "Sources: " + " · ".join(uniq), 10, color(theme, "muted"), italic=True)


def render_stat(slide, plan, theme, s):
    fill_bg(slide, theme, "bg_light")
    tf = add_box(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
    style_run(tf.paragraphs[0].add_run(), s.get("title", ""), 28, color(theme, "primary"), bold=True)

    stat = s.get("stat", {}) or {}
    is_review = stat.get("status") == "needs_review"
    val_color = color(theme, "warn") if is_review else color(theme, "text_dark")

    vtf = add_box(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
    vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    style_run(vtf.paragraphs[0].add_run(), str(stat.get("value", "")), 66, val_color, bold=True)

    ltf = add_box(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0))
    ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
    style_run(ltf.paragraphs[0].add_run(), stat.get("label", ""), 18, color(theme, "muted"))
    if is_review:
        m = ltf.add_paragraph()
        m.alignment = PP_ALIGN.CENTER
        style_run(m.add_run(), "⚠ needs review", 13, color(theme, "warn"), bold=True)
    elif stat.get("status") == "verified":
        lbl = source_label(plan, stat.get("source"))
        if lbl:
            foot = add_box(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
            foot.paragraphs[0].alignment = PP_ALIGN.CENTER
            style_run(foot.paragraphs[0].add_run(), "Source: " + lbl, 10, color(theme, "muted"), italic=True)


def render_to_verify(slide, plan, theme, items):
    fill_bg(slide, theme, "bg_dark")
    tf = add_box(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
    style_run(tf.paragraphs[0].add_run(), "To verify before presenting", 30, color(theme, "warn"), bold=True)

    body = add_box(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.0))
    if not items:
        style_run(body.paragraphs[0].add_run(),
                  "Nothing outstanding — every claim is verified or owned by you. ✅",
                  18, color(theme, "text_light"))
        return
    first = True
    for sidx, text in items:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.space_after = Pt(10)
        style_run(p.add_run(), f"⚠  ", 16, color(theme, "warn"), bold=True)
        style_run(p.add_run(), f"(slide {sidx}) ", 14, color(theme, "muted"))
        style_run(p.add_run(), text, 16, color(theme, "text_light"))


RENDERERS = {
    "title": render_title,
    "content": render_content,
    "stat": render_stat,
    "closing": render_closing,
}


def main():
    ap = argparse.ArgumentParser(description="Render a deck plan to .pptx")
    ap.add_argument("deck_plan")
    ap.add_argument("out_pptx")
    args = ap.parse_args()

    with open(args.deck_plan, "r", encoding="utf-8") as f:
        plan = json.load(f)

    theme = dict(DEFAULT_THEME)
    theme.update(plan.get("theme", {}) or {})

    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    review_items = []
    for s in plan.get("slides", []):
        slide = prs.slides.add_slide(blank)
        layout = s.get("layout", "content")
        RENDERERS.get(layout, render_content)(slide, plan, theme, s)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
        # Collect needs_review claims for the auto closing slide.
        sidx = len(prs.slides._sldIdLst)  # 1-based number of the slide just added
        for b in s.get("bullets", []) or []:
            if b.get("status") == "needs_review":
                review_items.append((sidx, b.get("text", "")))
        stat = s.get("stat")
        if isinstance(stat, dict) and stat.get("status") == "needs_review":
            review_items.append((sidx, f"{stat.get('value', '')} — {stat.get('label', '')}".strip(" —")))

    # Auto "To Verify" slide.
    slide = prs.slides.add_slide(blank)
    render_to_verify(slide, plan, theme, review_items)

    prs.save(args.out_pptx)
    print(f"Wrote {args.out_pptx} — {len(prs.slides._sldIdLst)} slides "
          f"({len(review_items)} need review).")
    return 0


if __name__ == "__main__":
    sys.exit(main())
